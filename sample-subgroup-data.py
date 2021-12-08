import os
import json
import re
from collections import defaultdict

import pptx.shapes.group
from tqdm import tqdm
import random
import copy
import numpy as np

from pptx import Presentation

# if two shape overlap area >= OVERLAP_MAX, they are considered as inside/outside
OVERLAP_MAX = 0.7

edge_dict2 = {2: "downside", 3: "lower right", 4: "right",
              5: "upper right", 6: "inside"}

vertex_dict = {"vertex_dict_flag": 2}
shape_appear_dict = {}
shape_ignore_set = set()
graph_num_slide_dict = {}

def get_vertex_label(shape_name):
    # save the shape type in vertex_dict & return the idx of current shape type
    norm_shape_name = shape_name.strip().replace(" ", "").lower()
    if norm_shape_name not in vertex_dict.keys():
        vertex_dict[norm_shape_name] = vertex_dict["vertex_dict_flag"]
        vertex_dict["vertex_dict_flag"] += 1
    return vertex_dict[norm_shape_name]

def cal_overlap(shape_bound1, shape_bound2):
    x1 = max(shape_bound1[0], shape_bound2[0])
    x2 = min(shape_bound1[0] + shape_bound1[2], shape_bound2[0] + shape_bound2[2])

    y1 = max(shape_bound1[1], shape_bound2[1])
    y2 = min(shape_bound1[1] + shape_bound1[3], shape_bound2[1] + shape_bound2[3])

    square1 = (x2 - x1) * (y2 - y1)
    square2 = shape_bound2[2] * shape_bound2[3]

    if square2 == 0:
        return False

    if float(square1) / square2 >= OVERLAP_MAX:
        return True
    else:
        return False

def get_edge_label(shape_bound1, shape_bound2):
    '''
    change for direction subgraph mining, change 10 relation into 5
    :param shape_bound1:[x1,y1,w1,h1]
    :param shape_bound2:[x2,y2,w2,h2]
    :return: boolean -- flag whether need change shape1 & shape2 bition
             int -- means the bition relation between shapes
    '''
    x1_l = shape_bound1[0]
    x1_r = shape_bound1[0] + shape_bound1[2]
    x1 = float(shape_bound1[0] * 2 + shape_bound1[2]) / 2.0

    x2_l = shape_bound2[0]
    x2_r = shape_bound2[0] + shape_bound2[2]
    x2 = float(shape_bound2[0] * 2 + shape_bound2[2]) / 2.0

    y1_u = shape_bound1[1]
    y1_d = shape_bound1[1] + shape_bound1[3]
    y1 = float(shape_bound1[1] * 2 + shape_bound1[3]) / 2.0

    y2_u = shape_bound2[1]
    y2_d = shape_bound2[1] + shape_bound2[3]
    y2 = float(shape_bound2[1] * 2 + shape_bound2[3]) / 2.0

    if abs(x1 - x2) >= 0.8 or abs(y1 - y2) >= 0.8:
        return False, -1

    if x2 < x1_l:
        if y2 < y1_u:
            return True, 3
        elif y2 <= y1_d:
            return True, 4
        else:
            return True, 5
    elif x2 <= x1_r:
        if y2 < y1_u:
            return True, 2
        elif y2 <= y1_d:
            if x1_l <= x2_l and x2_r <= x1_r and y1_u <= y2_u and y2_d <= y1_d:
                return False, 6
            elif x2_l <= x1_l and x1_r <= x2_r and y2_u <= y1_u and y1_d <= y2_d:
                return True, 6
            elif cal_overlap(shape_bound1, shape_bound2):
                return False, 6
            elif cal_overlap(shape_bound2, shape_bound1):
                return True, 6
            else:
                if x2 < x1:
                    if y2 < y1:
                        return True, 3
                    elif y2 == y1:
                        return True, 4
                    else:
                        return True, 5
                elif x2 == x1:
                    if y2 < y1:
                        return True, 2
                    elif y2 == y1:
                        print("????????")
                        return False, -1
                    else:
                        return False, 2
                else:
                    if y2 < y1:
                        return False, 5
                    elif y2 == y1:
                        return False, 4
                    else:
                        return False, 3
        else:
            return False, 2
    else:
        if y2 < y1_u:
            return False, 5
        elif y2 <= y1_d:
            return False, 4
        else:
            return False, 3

def add_appear_time(shape_label):
    if shape_label not in shape_appear_dict.keys():
        shape_appear_dict[shape_label] = 0
    shape_appear_dict[shape_label] += 1

def find_origin_shape(shape_idx,shape_list):
    for idx,shape_tmp in enumerate(shape_list):
        if str(shape_tmp["ShapeIdx"]).strip() == shape_idx:
            return shape_tmp
    return None

class OwnShape:
    def __init__(self,shape_idx,shape_name_strip,bound):
        self.shape_idx = shape_idx
        self.shape_name = shape_name_strip
        self.bound = bound
        # self.x = bound["R_Left"]
        # self.y = bound["R_Top"]
        # self.w = bound["R_Width"]
        # self.h = bound["R_Height"]

    def Bound(self):
        #return [self.x,self.y,self.w,self.h]
        return self.bound

    def shapeIdx(self):
        return self.shape_idx

    def shapeName(self):
        return self.shape_name


class OwnGroup:
    def __init__(self,file_idx,page_idx,shape_idx):
        self.file_idx = file_idx
        self.page_idx = page_idx
        self.ori_shape_idx = shape_idx
        self.shapes = []
        self.shape_info_set = set()
        self.shape_name_set = set()

    def add_ownshape(self,ownshape):
        self.shapes.append(ownshape)
        self.shape_info_set.add(ownshape.shape_idx)
        self.shape_name_set.add(ownshape.shape_name)

    def add_shapes(self,shape_list):
        for iidx,shapp in enumerate(shape_list):
            #new_shape = OwnShape(shapp["ShapeIdx"].strip(), shape_name_strip(shapp), shapp["Bound"])
            new_shape = OwnShape(str(shapp["ShapeIdx"]).strip(), shapp["ShapeName"], shapp["Bound"])
            self.shapes.append(new_shape)
            self.shape_info_set.add(new_shape.shape_idx)
            self.shape_name_set.add(new_shape.shape_name)

    def add_ori_group_info(self,shape_name_strip,bound):
        self.ori_shape_name = shape_name_strip # 没有strip过！！！
        self.shape_info_set.add(self.ori_shape_idx)
        self.shape_name_set.add(self.ori_shape_name)
        self.ori_bound = bound  # dict

    def clear_shapes(self):
        self.clear_shape_list = []
        group_idx_name_set = set()
        for idx,current_shape in enumerate(self.shapes):
            if current_shape.shape_idx.find(".0") != -1 and current_shape.shape_idx.split(".0")[-1] == "":
               group_idx_name_set.add(current_shape.shape_idx[:-2])
        for idx,current_shape in enumerate(self.shapes):
            if current_shape.shape_idx not in group_idx_name_set:
                self.clear_shape_list.append(current_shape)

        return self.clear_shape_list


def get_shape_absolute_bound(tmp_group,filePath_slide,filePath_ori,shapeName_idx_dict):
    '''
    get the shape absolute bound in the group
    :param filePath_slide: the pptx files path
    :param filePath_ori: the json files path
    :param shapeName_idx_dict: the dict from shapeName to shapeIdx
    :return: a dict which keys are shapeIdx, values are shape bounds
    '''
    def get_bound(tmp_shape):
        if tmp_shape.name in tmp_group.shape_name_set:
            group_shape_bound_dict[shapeName_idx_dict[tmp_shape.name]] = [float(tmp_shape.left) / prs.slide_width,
                                                                      float(tmp_shape.top) / prs.slide_height,
                                                                      float(tmp_shape.width) / prs.slide_width,
                                                                      float(tmp_shape.height) / prs.slide_height,
                                                                      tmp_shape.rotation]
            tmp_set.add(tmp_shape.name)

    def recursion_groupshape(tmp_shape):
        tmp_shape_list = tmp_shape.shapes
        for shape_i in tmp_shape_list:
            get_bound(shape_i)
            if type(shape_i) == pptx.shapes.group.GroupShape:
                recursion_groupshape(shape_i)


    # shape name set
    src_dir = filePath_slide
    fileName_list = os.listdir(filePath_ori)
    filename = fileName_list[tmp_group.file_idx]

    group_shape_bound_dict = {}
    tmp_set = set()

    try:
        prs = Presentation(os.path.join(src_dir, filename[:-5]))
        prs_slide = prs.slides[tmp_group.page_idx]
        for shape in prs_slide.shapes:
           get_bound(shape)
           if type(shape) == pptx.shapes.group.GroupShape:
                recursion_groupshape(shape)
    except Exception as e:
        print(e)

    if len(tmp_set) != len(tmp_group.shape_name_set):
        print("OH NO!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!!")
        print(len(tmp_set),len(tmp_group.shape_name_set))
        print(tmp_set)
        print(tmp_group.shape_name_set)

    return group_shape_bound_dict

def merge_two_dict(x,y):
    z = x.copy()
    z.update(y)
    return z

def get_group_info(slide,file_idx,slide_idx):
    '''
    - extract the group info stored in the json files
    P.S. the group info is stored in the shape_idx of a shape
        e.g. "5.0","5.1","5.2"... is subgroup items of "5"
             "5.0.0" is a subgroup item of "5.0"

    - use class "OwnShape" and "OwnGroup" to store info and automatically

    :return:
        (both use shape_idx to represent a shape)
        results: a list of all groups in the layout, the last element in the list is nogroup_shape
            e.g. Group_real_list for the Picture "A-simple-example.png" in the github folder:
                [["3.1.1","3.1.0","3.1"],["3.1","3.0","3"],["3","2","1"]]
                - if shapeIdx is too hard to read, you can consider the letter below as shapeIdx to understand easier:
                  [[a,b,G1],[G1,c,G2],[G2,d,e]]

        group_dict: a dict contains all groups and the shapes in each group
            e.g. Group_real_list for the Picture "A-simple-example.png" in the github folder:
                {0:{"3":["3","3.0","3.1"]},1:{"3.1":["3.1","3.1.1","3.1.0"]}}
                - if shapeIdx is too hard to read, you can consider the letter below as shapeIdx to understand easier:
                  {0:{G2:[G2,c,G1]},1:{G1:[G1,a,b]}}
    '''
    # a dict to store all group node's json info in a group
    shape_idx_dict = defaultdict(list)

    shapeName_idx_dict = {}

    # single nodes which has no group & the group nodes which on the same level of single nodes
    ## [shape_idx1,shape_idx2,......]
    nogroup_shape = []

    # a list contains all groups in type of Owngroup
    tmp_group_list = []

    # a dict contains all groups and the shapes in each group
    ## 0 for groups on the same level of single shape in layout
    ## 1 for groups which is the subgroup of groups-0
    group_dict = {0:{},1:{}}

    # save each shape into corresponding group
    for shape in slide["OrderedShapes"]:
        try:
            shape_idx = str(shape["ShapeIdx"]).strip()
            shapeName_idx_dict[shape["ShapeName"]] = shape_idx

            shape_idx_dict[shape_idx].append(shape)
            if shape_idx.find(".") == -1:
                # new_shape = OwnShape(str(shape["ShapeIdx"]).strip(), shape["ShapeName"], shape["Bound"])
                nogroup_shape.append(str(shape["ShapeIdx"]).strip())

            if shape_idx.find(".") != -1: # the shape is a element of a group

                flag_shape_idx = shape_idx
                shape_idx_tmp = flag_shape_idx.split(".")[-1]
                tmp_group_idx = flag_shape_idx[:-len(shape_idx_tmp) - 1]
                shape_idx_dict[tmp_group_idx].append(shape)

                # find the group which on the sample level of single nodes
                if shape_idx.find(".0") != -1: # for the first element of a group
                    shape_idxs_list = shape_idx.split(".0")

                    if len(shape_idxs_list) == 2 and shape_idxs_list[1] == "" and shape_idxs_list[0].find(".") == -1:
                        group_dict[0][shape_idxs_list[0]] = "Meaningless"

        except Exception as e:
            print("Error contains in ", file_idx, slide_idx, shape_idx, shape["ShapeIdx"])
            print(e)

    # get all groups and their elements in layout
    for shape2 in slide["OrderedShapes"]:
        try:
            shape_idx = str(shape2["ShapeIdx"]).strip()

            if shape_idx.find(".") != -1:
                if shape_idx.find(".0") != -1 and shape_idx.split(".0")[-1] == "":
                    new_group = OwnGroup(file_idx, slide_idx, shape_idx[:-2])
                    tmp_group_list.append(new_group)

                    ori_group_shape_list = shape_idx_dict[shape_idx[:-2]]
                    ori_group_shape = find_origin_shape(shape_idx[:-2], ori_group_shape_list)
                    if ori_group_shape == None:
                        print("?????")
                    # group_list[-1].add_ori_group_info(shape_name_strip(ori_group_shape),ori_group_shape["Bound"])
                    tmp_group_list[-1].add_ori_group_info(ori_group_shape["ShapeName"],
                                                               ori_group_shape["Bound"])
                    ori_group_shape_list.remove(ori_group_shape)
                    tmp_group_list[-1].add_shapes(ori_group_shape_list)

                    tmp_group_list[-1].clear_shapes()

                    if tmp_group_list[-1].ori_shape_idx not in group_dict[0]:
                        group_dict[1][tmp_group_list[-1].ori_shape_idx] =list(tmp_group_list[-1].shape_info_set)
                    else:
                        group_dict[0][tmp_group_list[-1].ori_shape_idx] =list(tmp_group_list[-1].shape_info_set)

        except Exception as e:
            print("Error contains in ", file_idx, slide_idx, shape_idx)
            print(e)

    # a dict contains all shape and their absolute bound
    shapeIdx_Bound_dict = {}

    results = []
    for tmp_group in tmp_group_list:
        results.append(list(tmp_group.shape_info_set))
        tmp_bound_dict = get_shape_absolute_bound(tmp_group,r"./pptx-path",r"./json-result",shapeName_idx_dict)
        shapeIdx_Bound_dict = merge_two_dict(shapeIdx_Bound_dict,tmp_bound_dict)
    results.append(nogroup_shape)

    return results,group_dict,shapeIdx_Bound_dict


def all_subgraphs(origin_list,min_size,max_size):
    # return all subgraph of origin_list which size between min_size and max_size
    # sigma(C^k_n)
    if len(origin_list) > 5:
        origin_list = random.sample(origin_list,5)
    results = [origin_list]
    full_results = [[origin_list]]
    tmp_length = len(origin_list)
    while (tmp_length > min_size):
        tmp_list = full_results[-1]
        tmp_result_list = []
        for list_i in tmp_list:
            for node_i in list_i:
                tmp = copy.deepcopy(list_i)
                tmp.remove(node_i)
                tmp_result_list.append(tmp)
        tmp_result_list = [list(t) for t in set(tuple(_) for _ in tmp_result_list)]
        full_results.append(tmp_result_list)
        results += tmp_result_list
        tmp_length -= 1

    return results

def get_group_shapetype(group1,shape_dict,group_full_dict):
    type_set1 = set()
    if group1 in group_full_dict[0].keys():
        id = 0
    else:
        id = 1
    for shape_i in group_full_dict[id][group1]:
        shape = shape_dict[shape_i]
        shape_name = shape["ShapeName"]
        tmp_shape_id = re.findall(r"\d+", shape_name)
        if len(tmp_shape_id) != 0:
            tmp_shape_name_id = int(tmp_shape_id[0])
            shape_name_strip = shape_name.replace(str(tmp_shape_name_id), "").strip()
        else:
            shape_name_strip = shape_name.strip()
        shape_label = get_vertex_label(shape_name_strip)

        type_set1.add(shape_label)

    return type_set1

def generate_same_type_groups_sample(groups,shape_dict,group_full_dict):
    results = {}
    for group_i in groups:
        tmp_set = tuple(get_group_shapetype(group_i,shape_dict,group_full_dict))
        if tmp_set not in results.keys():
            results[tmp_set] = []
        results[tmp_set].append(group_i)

    sample_results = []
    for key_i in results.keys():
        if len(results[key_i]) > 1:
            sample_results += all_subgraphs(results[key_i],2,len(results[key_i]))
    return sample_results

def subgroup_data_generation(filePath, outputPath, tmp_s,full_edge=False,front_mode=True,group_mode=True,stop_number=100):

    # filePath: the json files path of vault-2300 dataset
    fileName_list = os.listdir(filePath)

    # save the order of fileName_list to better visualization
    with open("./subgroup-data/fileName-list" + tmp_s, "w", encoding="utf-8")as f:
        for file_idx, fileName in enumerate(fileName_list):
            f.write(str(file_idx) + ": ")
            f.write(fileName + "\n")

    '''
    DATA PREPARATION - 1/3
    : get shape_ignore_set
    '''

    # get shape ignore set (ignore the shape type which appear time <= 50 in the whole vault-2300 data)
    for file_idx in tqdm(range(len(fileName_list))):

        fileName = fileName_list[file_idx]
        with open(os.path.join(filePath, fileName), "rb")as f:
            fileResult = json.load(f)
        slide_list = fileResult["Slides"]

        for slide_idx, slide in enumerate(slide_list):
            for shape_idx, shape in enumerate(slide["OrderedShapes"]):
                try:
                    # "strip" the shape name: "Group 5" -> "group"
                    shape_name = shape["ShapeName"]
                    tmp_shape_id = re.findall(r"\d+", shape_name)

                    if len(tmp_shape_id) != 0:
                        tmp_shape_name_id = int(tmp_shape_id[0])
                        shape_name_strip = shape_name.replace(str(tmp_shape_name_id), "").strip()
                    else:
                        shape_name_strip = shape_name.strip()

                    shape_label = get_vertex_label(shape_name_strip)
                    add_appear_time(shape_label)

                except Exception as e:
                    print("Error contains in ", file_idx, slide_idx, shape_idx, shape["ShapeIdx"])
                    print(e)

    for shape_label_t in shape_appear_dict:
        if shape_appear_dict[shape_label_t] <= 50:
            shape_ignore_set.add(shape_label_t)

    # save the shape_ignore_set, shape_appear_dict, shape_name_dict
    with open("./subgroup-data/shape-ignore" + tmp_s, "w", encoding="utf-8")as f:
        f.write(str(shape_ignore_set))
    with open("./subgroup-data/shape-appear" + tmp_s, "w", encoding="utf-8")as f:
        f.write(str(shape_appear_dict))
    with open("./subgroup-data/shape-name-dict" + tmp_s, "w", encoding="utf-8")as f:
        f.write(str(vertex_dict))

    '''
    DATA PREPARATION - 2/3
    : get group-related info in the layout
    '''

    graph_num = 0
    flag = True # flag for control the number of generated data

    graph_num_slide_strings = []

    # save the string type of sampled data
    write_pos_a_strings = [] #pos_target_graph
    write_pos_b_strings = [] #pos_query_graph
    write_neg_a_strings = [] #neg_target_graph
    write_neg_b_strings = [] #neg_query_graph

    #sum the string list
    overall_data = [write_pos_b_strings,write_neg_b_strings,write_pos_a_strings,write_neg_a_strings]

    if front_mode == True:
        range_list = range(len(fileName_list))
    else:
        range_list = range(len(fileName_list)-1,-1,-1)
    for file_idx in tqdm(range_list):
        fileName = fileName_list[file_idx]

        if flag == False:
            break

        '''
        The dataset vault-200(contains 200 grouped graph,we use to train the model before) just use ppt-json-file 0~53
        So vault-200 and the "subgroup dataset" are sampled from the same graph data, only different in sample method
        '''
        if file_idx >= 54:
            print("file enough")
            break

        with open(os.path.join(filePath, fileName), "rb")as f:
            fileResult = json.load(f)
        slide_list = fileResult["Slides"]
        for slide_idx, slide in enumerate(slide_list):
            # 2 < shape_number < 100
            if len(slide["OrderedShapes"]) <= 2 or len(slide["OrderedShapes"]) >= 100:
                continue

            if graph_num == stop_number:
                print("Enough graph!")
                flag = False
                break

            # get group-related info in the slide
            group_shapeIdx_list,group_dict,shapeIdx_Bound_dict = get_group_info(slide, file_idx, slide_idx)

            if group_mode:# ignore the slide which has no group
                if len(group_shapeIdx_list) == 1:  # group data only # group-1
                    continue
            else:
                if len(group_shapeIdx_list) == 1 and len(slide["OrderedShapes"]) >= 30:  # simple slides # group-0
                    continue

            '''
            DATA PREPARATION - 3/3
            : delete the ignore shape type in group_shapeIdx_list,group_dict
            '''

            # dict for transform shapeIdx into shape: {shapeIdx1 : shape json1,shapeIdx2 : shape json2......}
            shapeIdx_shape_dict = {}

            # find the ignore shape in this slide using shapes' type, save ignore shapeIdx in tmp_shape_ignore_set
            tmp_shape_ignore_set = set()
            for shape_idx, shape in enumerate(slide["OrderedShapes"]):
                try:
                    # get shape name
                    shape_name = shape["ShapeName"]
                    tmp_shape_id = re.findall(r"\d+", shape_name)
                    shapeIdx_shape_dict[shape["ShapeIdx"]] = shape

                    if len(tmp_shape_id) != 0:
                        tmp_shape_name_id = int(tmp_shape_id[0])
                        shape_name_strip = shape_name.replace(str(tmp_shape_name_id), "").strip()
                    else:
                        shape_name_strip = shape_name.strip()

                    # judge the current shape should be ignore or not
                    shape_label = get_vertex_label(shape_name_strip)
                    if shape_label in shape_ignore_set:
                        tmp_shape_ignore_set.add(shape["ShapeIdx"])
                        continue
                except Exception as e:
                    print("Error contains in ", file_idx, slide_idx, shape_idx, shape["ShapeIdx"])
                    print(e)

            # a list contains all groups in the slide without ignore shape type
            # e.g. Group_real_list for the Picture "A-simple-example.png" in the github folder:(consider the letters as the shapeIdx)
            #     [[a,b,G1],[G1,c,G2],[G2,d,e]]
            group_real_list = []
            for tmp_group in group_shapeIdx_list:
                tmp_group_list = []
                for tmp_shape_idx in tmp_group:
                    if tmp_shape_idx not in tmp_shape_ignore_set:
                        tmp_group_list.append(tmp_shape_idx)
                if len(tmp_group_list) != 0:
                    group_real_list.append(tmp_group_list)

            # a dict contains all groups
            # e.g. Group_real_list for the Picture "A-simple-example.png" in the github folder:(consider the letters as the shapeIdx)
            #     {0:{G2:[G2,c,G1]},1:{G1:[G1,a,b]}}
            group_full_dict = {0:{},1:{}}
            for group_ii in group_dict[0].keys():
                if group_ii not in tmp_shape_ignore_set:
                    group_full_dict[0][group_ii] = group_dict[0][group_ii]
            for group_jj in group_dict[1].keys():
                if group_jj not in tmp_shape_ignore_set:
                    group_full_dict[1][group_jj] = group_dict[1][group_jj]

            '''
            DATA GENERATION - POS
            '''

            # start pos data generation
            pos_a_data_list = []
            pos_b_data_list = []

            '''
            pos-data-type1 : 1/2

            - Range: for all groups
            - Query_graph(pos_b): some nodes in this group
            - Target_graph(pos_a): all nodes in this group
            '''
            for id in range(1):
                for group_i in group_full_dict[id].keys():
                    # single node ~ all nodes in the graph
                    # e.g. <G1,c> -> <G1,c,G2>
                    results_list = all_subgraphs(group_full_dict[id][group_i],2,len(group_full_dict[id][group_i]))

                    # control the number of results_list
                    # the length of a 4-item list's all_subgraphs function result == 11
                    if len(results_list) > 11:
                        results_list = random.sample(results_list,11)
                    for i in range(len(results_list)):
                        pos_a_data_list.append(list(set(group_full_dict[id][group_i]))) # set() for remove duplication
                        pos_b_data_list.append(list(set(results_list[i])))

                    # for nested group
                    small_group = []
                    for group_j in group_full_dict[id][group_i]:
                        if group_j in group_full_dict[1].keys():
                            small_group.append(group_j)

                            # e.g. <a,b> -> <G1,c,G2,a,b>
                            results_list = all_subgraphs(group_full_dict[1][group_j], 2,
                                                         len(group_full_dict[1][group_j]))
                            if len(results_list) > 11:
                                results_list = random.sample(results_list, 11)
                            for i in range(len(results_list)):
                                pos_a_data_list.append(list(set(group_full_dict[id][group_i] + group_full_dict[1][group_j])))  # 自身
                                pos_b_data_list.append(list(set(results_list[i])))

                    '''
                    pos-data-type2 : 1/2

                    - Range: for "small" groups
                    - Query_graph(pos_b): all nodes in one group or all groups
                    - Target_graph(pos_a): all nodes in all groups (+ big group items)
                    '''
                    if len(small_group) > 1:
                        # two or more small(nested) groups which have sample set of shape type
                        group_results = generate_same_type_groups_sample(small_group,shapeIdx_shape_dict,group_full_dict)
                        if len(group_results) > 20:
                            group_results = random.sample(group_results, 20)
                        for tmp_i in range(len(group_results)):
                            # num = random.randint(2,len(small_group))
                            tmp_sample_list = group_results[tmp_i]
                            tmp_a_data = group_full_dict[id][group_i]
                            tmp_b_data = []
                            for group_iii in tmp_sample_list:
                                tmp_a_data += group_full_dict[1][group_iii]
                                tmp_b_data += group_full_dict[1][group_iii]
                            # <a-1,b-1,G1-1,a-2,b-2,G1-2> -> <a-1,b-1,G1-1,a-2,b-2,G1-2,G2,c>
                            # picture example-for-line639-660.png
                            pos_a_data_list.append(list(set(tmp_a_data)))
                            pos_b_data_list.append(list(set(tmp_b_data)))

                            # <a-1,b-1,G1-1> -> <a-1,b-1,G1-1,a-2,b-2,G1-2>
                            # picture example-for-line639-660.png
                            for group_jjj in tmp_sample_list:
                                pos_a_data_list.append(list(set(tmp_b_data)))
                                pos_b_data_list.append(list(set(group_full_dict[1][group_jjj])))



            '''
            pos-data-type2 : 2/2

            - Range: for "big" groups
            - Query_graph(pos_b): all nodes in one group or
            - Target_graph(pos_a): all nodes in this group
            '''
            if len(group_real_list[-1]) -1 >= 1:

                tmp_single_node_list = copy.deepcopy(group_real_list[-1])
                for group_j in group_full_dict[0].keys():
                    if group_j in tmp_single_node_list:
                        tmp_single_node_list.remove(group_j)

                big_group_num  = len(group_full_dict[0].keys())
                if big_group_num > 1:
                    group_results = generate_same_type_groups_sample(group_full_dict[0].keys(), shapeIdx_shape_dict,group_full_dict)

                    if len(group_results) > 20:
                        group_results = random.sample(group_results,20)

                    for tmp_i in range(len(group_results)):
                        sample_big_groups = group_results[tmp_i]

                        sample_single_group = sample_big_groups[0]
                        tmp_a_data = copy.deepcopy(group_real_list[-1])
                        tmp_b_data = []
                        for group_iii in sample_big_groups:
                            if group_iii != sample_single_group:
                                tmp_a_data += group_full_dict[0][group_iii]
                                tmp_b_data += group_full_dict[0][group_iii]
                        # entire layout : all nodes in groups
                        pos_a_data_list.append(list(set(tmp_a_data + group_full_dict[0][sample_single_group])))
                        pos_b_data_list.append(list(set(tmp_b_data + group_full_dict[0][sample_single_group])))

                        for group_jj in sample_big_groups:  # all nodes in groups : nodes in one groups
                            pos_a_data_list.append(list(set(tmp_b_data + group_full_dict[0][sample_single_group])))
                            pos_b_data_list.append(list(set(group_full_dict[0][group_jj])))

                        results_list = all_subgraphs(group_full_dict[0][sample_single_group], 2,
                                                     len(group_full_dict[0][sample_single_group]))
                        if len(results_list) > 11:
                            results_list = random.sample(results_list, 11)
                        for i in range(len(results_list)):
                            pos_a_data_list.append(
                                list(set(tmp_b_data + group_full_dict[0][sample_single_group])))  # <a,b> & <G2...>
                            pos_b_data_list.append(list(set(results_list[i])))

                            pos_a_data_list.append(list(set(tmp_a_data + group_full_dict[0][sample_single_group])))
                            pos_b_data_list.append(list(set(results_list[i])))

                            pos_a_data_list.append(
                                list(set(tmp_b_data + group_full_dict[0][sample_single_group])))  # 满足一个组后，才考虑其他零碎点？
                            pos_b_data_list.append(list(set(results_list[i] + tmp_b_data)))

                            pos_a_data_list.append(list(set(tmp_a_data + group_full_dict[0][sample_single_group])))
                            pos_b_data_list.append(list(set(results_list[i] + tmp_b_data)))

                        # <G1,c,G2,d> is subgraph of <G1,c,G2,d,e>
                        if len(tmp_single_node_list) <= 5:
                            results_list2 = all_subgraphs(tmp_single_node_list, 2,
                                                          len(tmp_single_node_list))
                            if len(results_list2) > 11:
                                results_list2 = random.sample(results_list2, 11)
                            for i in range(len(results_list2)):
                                pos_a_data_list.append(
                                    list(set(tmp_a_data + group_full_dict[0][sample_single_group])))
                                pos_b_data_list.append(
                                    list(set(tmp_b_data + group_full_dict[0][sample_single_group] + results_list2[i])))

                else:
                    '''
                    pos-data-type1 : 2/2

                    - Range: for single nodes and "big" groups
                    - Query_graph(pos_b): some nodes in this group
                    - Target_graph(pos_a): all nodes in this group
                    '''
                    for group_j in group_full_dict[0].keys():
                        pos_a_data_list.append(list(set(group_real_list[-1] + group_full_dict[0][group_j])))  # entire layout
                        pos_b_data_list.append(list(set(group_full_dict[0][group_j])))

                        # <G1,c> is subgraph of <G1,c,G2,d,e>
                        results_list = all_subgraphs(group_full_dict[0][group_j], 2,
                                                     len(group_full_dict[0][group_j]))
                        if len(results_list) > 11:
                            results_list = random.sample(results_list, 11)
                        for i in range(len(results_list)):
                            pos_a_data_list.append(
                                list(set(group_real_list[-1] + group_full_dict[0][group_j])))
                            pos_b_data_list.append(list(set(results_list[i])))

                        # <G1,c,G2,d> is subgraph of <G1,c,G2,d,e>
                        if len(tmp_single_node_list) <= 5:
                            results_list2 = all_subgraphs(tmp_single_node_list, 2,
                                                          len(tmp_single_node_list))
                            if len(results_list2) > 11:
                                results_list2 = random.sample(results_list2, 11)
                            for i in range(len(results_list2)):
                                pos_a_data_list.append(
                                    list(set(group_real_list[-1] + group_full_dict[0][group_j])))
                                pos_b_data_list.append(list(set(results_list2[i] + group_full_dict[0][group_j])))

            '''
            DATA GENERATION - NEG
            '''
            # start neg data generation
            neg_a_data_list = []
            neg_b_data_list = []

            '''
            neg-data-type1 : 1/2

            - Range: for all groups
            - Query_graph(pos_b): not all nodes in nested group and nodes in the "outside" group
            - Target_graph(pos_a): all nodes in this group
            '''
            for id in range(1):
                for group_i in group_full_dict[id].keys():
                    for group_j in group_full_dict[id][group_i]:
                        if group_j in group_full_dict[1].keys():
                            # small_group.append(group_j)
                            tmp_list_i = copy.deepcopy(group_full_dict[id][group_i])
                            tmp_list_i.remove(group_j)
                            tmp_list_i.remove(group_i)
                            if len(tmp_list_i) <=1:
                                continue

                            tmp_list_j = copy.deepcopy(group_full_dict[1][group_j])
                            tmp_list_j.remove(group_j)
                            if len(tmp_list_j) <=2:
                                continue

                            for i in range(len(tmp_list_i)+len(tmp_list_j) - 3):
                                sample1 = random.randint(1, len(tmp_list_i))
                                sample1_results = random.sample(tmp_list_i, sample1)
                                sample2 = random.randint(2, len(tmp_list_j) - 1)
                                sample2_results = random.sample(tmp_list_j, sample2)

                                neg_a_data_list.append(list(set(group_full_dict[id][group_i] + group_full_dict[1][group_j])))
                                neg_b_data_list.append(list(set(sample1_results + sample2_results)))

            '''
            neg-data-type1 : 2/2

            - Range: for "big" groups
            - Query_graph(pos_b): not all nodes in "big" group and single nodes
            - Target_graph(pos_a): entire layout
            '''
            for group_i in group_full_dict[0].keys():
                tmp_list_i = copy.deepcopy(group_full_dict[0][group_i])
                tmp_list_i.remove(group_i)
                if len(tmp_list_i) <= 1:
                    continue

                tmp_list_j = copy.deepcopy(group_real_list[-1])
                tmp_list_j.remove(group_i)
                if len(tmp_list_j) <= 2:
                    continue

                for i in range(min(5,len(tmp_list_i) + len(tmp_list_j) - 3)):
                    sample1 = random.randint(1, min(len(tmp_list_i),5))
                    sample1_results = random.sample(tmp_list_i, sample1)
                    sample2 = random.randint(2, min(len(tmp_list_j) - 1,5))
                    sample2_results = random.sample(tmp_list_j, sample2)

                    neg_a_data_list.append(
                        list(set(group_real_list[-1] + group_full_dict[0][group_i])))
                    neg_b_data_list.append(list(set(sample1_results + sample2_results)))

            '''
            DATA TRANSFORMATION
            - from shape_idx list into the format model can read
            P.S. data format:
                "t # N" means the Nth graph,
                "v M L" means that the Mth vertex in this graph has label L,
                "e P Q L" means that there is an edge connecting the Pth vertex with the Qth vertex. The edge has label L.
            '''

            # b_data = [pos_b_data_list,neg_b_data_list]
            # a_data = [pos_a_data_list,neg_a_data_list]
            all_data = [pos_b_data_list,neg_b_data_list,pos_a_data_list,neg_a_data_list]

            # ignore for have no edges in graph
            ignore_pos = set()
            ignore_neg = set()
            b_data_ignore = [ignore_pos,ignore_neg]

            # print("example sum:",len(pos_a_data_list),len(neg_a_data_list))
            # print("assert equal ",len(pos_a_data_list),len(pos_b_data_list),end=",")
            # print(len(neg_a_data_list),len(neg_b_data_list))

            for iidx,data_list in enumerate(all_data):
                for list_idx,nodes_list in enumerate(data_list):
                    if iidx >= 2:
                        if list_idx in b_data_ignore[iidx - 2]:
                            continue

                    nodes_IdxLabel_list = []
                    nodes_Flag_list = []
                    for vertex_idx,vertex_i in enumerate(nodes_list):
                        shape = shapeIdx_shape_dict[vertex_i]
                        shape_name = shape["ShapeName"]
                        tmp_shape_id = re.findall(r"\d+", shape_name)

                        if len(tmp_shape_id) != 0:
                            tmp_shape_name_id = int(tmp_shape_id[0])
                            shape_name_strip = shape_name.replace(str(tmp_shape_name_id), "").strip()
                        else:
                            shape_name_strip = shape_name.strip()

                        shape_label = get_vertex_label(shape_name_strip)

                        # prepare for if node unreachable
                        nodes_IdxLabel_list.append((shape["ShapeIdx"],shape_label))
                        nodes_Flag_list.append(False)

                    edge_IdxLabel_list = []

                    if full_edge == True:
                        for shape_idx_i in range(len(nodes_list)):
                            for shape_idx_j in range(shape_idx_i + 1, len(nodes_list)):
                                # edge_label = get_edge_label(shape_list[shape_i],shape_list[shape_j])
                                # if edge_label == -1:
                                #     continue
                                shape1 = shapeIdx_shape_dict[nodes_list[shape_idx_i]]
                                shape2 = shapeIdx_shape_dict[nodes_list[shape_idx_j]]
                                if nodes_list[shape_idx_i] in shapeIdx_Bound_dict.keys():
                                    bound1 = shapeIdx_Bound_dict[nodes_list[shape_idx_i]]
                                else:
                                    bound1 = [shape1["Bound"]["R_Left"], shape1["Bound"]["R_Top"],
                                              shape1["Bound"]["R_Width"],
                                              shape1["Bound"]["R_Height"]]
                                if nodes_list[shape_idx_j] in shapeIdx_Bound_dict.keys():
                                    bound2 = shapeIdx_Bound_dict[nodes_list[shape_idx_j]]
                                else:
                                    bound2 = [shape2["Bound"]["R_Left"], shape2["Bound"]["R_Top"],
                                              shape2["Bound"]["R_Width"],
                                              shape2["Bound"]["R_Height"]]
                                edge_flag, edge_label = get_edge_label(bound1, bound2)
                                if edge_label == -1:
                                    continue
                                if edge_flag:
                                    # write_list.append("e " + str(shape_idx_j) + " " + str(shape_idx_i) + " " + str(edge_label) + "\n")
                                    edge_IdxLabel_list.append((nodes_list[shape_idx_j],nodes_list[shape_idx_i],edge_label))
                                    if nodes_Flag_list[shape_idx_i] == False:
                                        nodes_Flag_list[shape_idx_i] = True
                                    if nodes_Flag_list[shape_idx_j] == False:
                                        nodes_Flag_list[shape_idx_j] = True

                                else:
                                    # write_list.append("e " + str(shape_idx_i) + " " + str(shape_idx_j) + " " + str(edge_label) + "\n")
                                    edge_IdxLabel_list.append(
                                        (nodes_list[shape_idx_i], nodes_list[shape_idx_j], edge_label))
                                    if nodes_Flag_list[shape_idx_i] == False:
                                        nodes_Flag_list[shape_idx_i] = True
                                    if nodes_Flag_list[shape_idx_j] == False:
                                        nodes_Flag_list[shape_idx_j] = True

                    else: # group_edge , actually we use
                        for shape_idx_i in range(len(nodes_list)):
                            for shape_idx_j in range(shape_idx_i + 1, len(nodes_list)):

                                tmp_i_list = []
                                for flag_i,group_i in enumerate(group_real_list):
                                    if nodes_list[shape_idx_i] in group_i:
                                        tmp_i_list.append(flag_i)
                                tmp_j_list = []
                                for flag_j,group_j in enumerate(group_real_list):
                                    if nodes_list[shape_idx_j] in group_j:
                                        tmp_j_list.append(flag_j)

                                x = [k for k in tmp_i_list if k in tmp_j_list]
                                if len(x) == 0:
                                    continue

                                shape1 = shapeIdx_shape_dict[nodes_list[shape_idx_i]]
                                shape2 = shapeIdx_shape_dict[nodes_list[shape_idx_j]]
                                if nodes_list[shape_idx_i] in shapeIdx_Bound_dict.keys():
                                    #print("yepe")
                                    bound1 = shapeIdx_Bound_dict[nodes_list[shape_idx_i]]
                                else:
                                    bound1 = [shape1["Bound"]["R_Left"], shape1["Bound"]["R_Top"],
                                              shape1["Bound"]["R_Width"],
                                              shape1["Bound"]["R_Height"]]
                                if nodes_list[shape_idx_j] in shapeIdx_Bound_dict.keys():
                                    #print("yepe")
                                    bound2 = shapeIdx_Bound_dict[nodes_list[shape_idx_j]]
                                else:
                                    bound2 = [shape2["Bound"]["R_Left"], shape2["Bound"]["R_Top"],
                                              shape2["Bound"]["R_Width"],
                                              shape2["Bound"]["R_Height"]]
                                edge_flag, edge_label = get_edge_label(bound1, bound2)
                                if edge_label == -1:
                                    continue
                                if edge_flag:
                                    # write_list.append("e " + str(shape_idx_j) + " " + str(shape_idx_i) + " " + str(edge_label) + "\n")
                                    edge_IdxLabel_list.append(
                                        (nodes_list[shape_idx_j], nodes_list[shape_idx_i], edge_label))
                                    if nodes_Flag_list[shape_idx_i] == False:
                                        nodes_Flag_list[shape_idx_i] = True
                                    if nodes_Flag_list[shape_idx_j] == False:
                                        nodes_Flag_list[shape_idx_j] = True
                                else:
                                    #write_list.append("e " + str(shape_idx_i) + " " + str(shape_idx_j) + " " + str(edge_label) + "\n")
                                    edge_IdxLabel_list.append(
                                        (nodes_list[shape_idx_i], nodes_list[shape_idx_j], edge_label))
                                    if nodes_Flag_list[shape_idx_i] == False:
                                        nodes_Flag_list[shape_idx_i] = True
                                    if nodes_Flag_list[shape_idx_j] == False:
                                        nodes_Flag_list[shape_idx_j] = True

                    write_list = []
                    # print("Start writing2")

                    shapeIdx_nodesIdx_dict = {}
                    shapeIdx_ignore_set = set()
                    real_nodes_IdxLabel_list = []
                    for flag_idx in range(len(nodes_Flag_list)):
                        # shapeIdx_nodesIdx_dict[nodes_list[flag_idx]] = len(real_nodes_IdxLabel_list)
                        # real_nodes_IdxLabel_list.append(nodes_IdxLabel_list[flag_idx])
                        if nodes_Flag_list[flag_idx] == False:
                            if len(nodes_list) != 1:
                                shapeIdx_ignore_set.add(nodes_list[flag_idx])
                            # if len(nodes_list) != 1:
                            #     print("TvT")
                        else:
                            shapeIdx_nodesIdx_dict[nodes_list[flag_idx]] = len(real_nodes_IdxLabel_list)
                            real_nodes_IdxLabel_list.append(nodes_IdxLabel_list[flag_idx])

                    if len(real_nodes_IdxLabel_list) == 0:
                        # print("????? no nodes in",iidx,list_idx)
                        if iidx <2:
                            b_data_ignore[iidx].add(list_idx)
                        continue

                    # write_list.append("t # " + str(len(overall_data[iidx])) + "\n")  ###############
                    for idx in range(len(real_nodes_IdxLabel_list)):
                        write_shapeIdx,write_shapeLabel = real_nodes_IdxLabel_list[idx]
                        write_list.append("v " + str(shapeIdx_nodesIdx_dict[write_shapeIdx]) + " " + str(write_shapeLabel) + "\n")

                    for edge_idx in range(len(edge_IdxLabel_list)):
                        shape_i,shape_j,write_edge_label = edge_IdxLabel_list[edge_idx]
                        if (shape_i not in shapeIdx_ignore_set) and (shape_j not in shapeIdx_ignore_set):
                            write_list.append(
                                "e " + str(shapeIdx_nodesIdx_dict[shape_i]) + " " + str(shapeIdx_nodesIdx_dict[shape_j]) + " " + str(write_edge_label) + "\n")

                    if iidx == 0:
                        graph_num_slide_dict[len(overall_data[iidx])] = (file_idx, slide_idx)
                    overall_data[iidx].append(write_list)

    '''
    neg-data-type2 : 1/1

    - Range: for whole dataset
    - Query_graph(pos_b): random
    - Target_graph(pos_a): random
    '''
    neg_a_random = random.sample(write_pos_a_strings,3000)
    neg_b_random = random.sample(write_neg_b_strings + write_neg_a_strings + write_pos_b_strings,3000)
    write_neg_a_strings += neg_a_random
    write_neg_b_strings += neg_b_random

    print("assert pos data len:",len(write_pos_a_strings),len(write_pos_b_strings))
    print("assert neg data len:",len(write_neg_a_strings),len(write_neg_b_strings))

    # shuffle dataset
    idx_range_pos = np.arange(len(write_pos_a_strings))
    np.random.shuffle(idx_range_pos)
    write_pos_a_strings = list(np.array(write_pos_a_strings)[idx_range_pos])
    write_pos_b_strings = list(np.array(write_pos_b_strings)[idx_range_pos])

    idx_range_neg = np.arange(len(write_neg_a_strings))
    np.random.shuffle(idx_range_neg)
    write_neg_a_strings = list(np.array(write_neg_a_strings)[idx_range_neg])
    write_neg_b_strings = list(np.array(write_neg_b_strings)[idx_range_neg])

    # save dataset
    with open(outputPath+"-pos-a", "w", encoding="utf-8") as f:
        for idx,strings in enumerate(write_pos_a_strings):
            f.write("t # " + str(idx) + "\n")
            for string in strings:
                f.write(string)
        f.write("t # -1")

    with open(outputPath+"-pos-b", "w", encoding="utf-8") as f:
        for idx,strings in enumerate(write_pos_b_strings):
            f.write("t # " + str(idx) + "\n")
            for string in strings:
                f.write(string)
        f.write("t # -1")

    with open(outputPath+"-neg-a", "w", encoding="utf-8") as f:
        for idx,strings in enumerate(write_neg_a_strings):
            f.write("t # " + str(idx) + "\n")
            for string in strings:
                f.write(string)
        f.write("t # -1")

    with open(outputPath+"-neg-b", "w", encoding="utf-8") as f:
        for idx,strings in enumerate(write_neg_b_strings):
            f.write("t # " + str(idx) + "\n")
            for string in strings:
                f.write(string)
        f.write("t # -1")

    # save graph number - file_idx & slide_idx dict
    with open("./subgroup-data/graph_dict" + tmp_s, "w", encoding="utf-8")as f:
        f.write(str(graph_num_slide_dict))

def transform_data(fileName,output_name):

    node_labels = []
    edges = []
    edges_labels = []
    graph_idx = []

    node_flag = 1
    tmp_flag = 0

    graphs = dict()
    with open(fileName, 'r') as f:
        lines = [line.strip() for line in f.readlines()]
        tgraph, graph_cnt = None, 0
        for i, line in enumerate(lines):
            if int(graph_cnt)>=4096:
                break
            cols = line.split(' ')
            if cols[0] == 't':
                node_flag += tmp_flag
                tmp_flag = 0

                if cols[-1] == '-1' or int(graph_cnt) >= 4096:
                    break
                graph_cnt =cols[2]

            elif cols[0] == 'v':
                # node_num = int(cols[-1]) + node_flag
                node_labels.append(cols[-1])#str(node_num))
                graph_idx.append(str(int(graph_cnt)+1))
                tmp_flag += 1
            elif cols[0] == 'e':
                node_num1 = int(cols[1]) + node_flag
                node_num2 = int(cols[2]) + node_flag
                edges.append(str(node_num1)+", "+str(node_num2))
                edges_labels.append(cols[-1])
                #tgraph.add_edge(AUTO_EDGE_ID, cols[1], cols[2], cols[3])
        # adapt to input files that do not end with 't # -1'

    with open("./subgroup-data/vault"+output_name+"_node_labels.txt","w",encoding='utf-8')as f:
        for i in range(len(node_labels)):
            f.write(str(node_labels[i]))
            f.write("\n")

    with open("./subgroup-data/vault"+output_name+"_A.txt","w",encoding='utf-8')as f:
        for i in range(len(edges)):
            f.write(str(edges[i]))
            f.write("\n")

    with open("./subgroup-data/vault"+output_name+"_edge_labels.txt","w",encoding='utf-8')as f:
        for i in range(len(edges_labels)):
            f.write(str(edges_labels[i]))
            f.write("\n")

    with open("./subgroup-data/vault"+output_name+"_graph_indicator.txt","w",encoding='utf-8')as f:
        for i in range(len(graph_idx)):
            f.write(str(graph_idx[i]))
            f.write("\n")

    with open("./subgroup-data/vault"+output_name+"_graph_labels.txt","w",encoding='utf-8')as f:
        for i in range(4096):
            f.write("1")
            f.write("\n")

if __name__ == "__main__":
    # origin_data_json_path, save_dataset_path, save_detailed_info_path,
    # front_mode: sample begin at the start/end of the data
    # group_mode: edges add between the shapes in the same group, set false will add edges between every shapes in the layout
    # stop_number: will stop when len(data) == stop_number
    subgroup_data_generation("./json-result",
                            "./subgroup-data/test",
                            "-test", front_mode=True, group_mode=True, stop_number=8096)

    transform_data("./subgroup-data/test-pos-a", "tryposA")
    transform_data("./subgroup-data/test-pos-b", "tryposB")
    transform_data("./subgroup-data/test-neg-a", "trynegA")
    transform_data("./subgroup-data/test-neg-b", "trynegB")
